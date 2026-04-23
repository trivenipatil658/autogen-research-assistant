from pathlib import Path  # Cross-platform file path handling
from pptx import Presentation  # python-pptx library for creating PowerPoint files
from pptx.util import Inches, Pt  # Utility classes for slide dimensions and font sizes
import re  # Regular expressions for parsing slide headings from text


def _parse_ppt_outline(text: str) -> list[dict]:
    """Extract slide title + bullets from writer output."""
    slides = []    # List to hold all parsed slide dictionaries
    current = None  # Tracks the slide currently being built

    for line in text.splitlines():  # Process the writer output line by line
        line = line.strip()  # Remove leading/trailing whitespace
        if not line:
            continue  # Skip blank lines

        # Detect slide headings — matches patterns like "Slide 1:", "## Slide 1", "**Slide 1**"
        if re.match(r"^(slide\s*\d+|##|\*\*slide)", line, re.IGNORECASE):
            if current:
                slides.append(current)  # Save the previous slide before starting a new one
            # Clean up the heading to extract just the title text
            title = re.sub(r"^[#*\s]*(slide\s*\d*:?\s*)?", "", line, flags=re.IGNORECASE).strip("* :")
            current = {"title": title or line, "bullets": []}  # Start a new slide dict

        elif current is not None and line.startswith(("-", "•", "*")):
            # Line is a bullet point — strip the bullet character and add to current slide
            current["bullets"].append(line.lstrip("-•* "))

        elif current is not None:
            # Line is plain text under a slide — treat it as a bullet point too
            current["bullets"].append(line)

    if current:
        slides.append(current)  # Don't forget to save the last slide

    return slides  # Return the list of slide dicts with title and bullets


def export_ppt(writer_text: str, topic: str, output_dir: Path) -> Path:
    prs = Presentation()  # Create a new blank PowerPoint presentation
    blank_layout = prs.slide_layouts[1]  # Use layout 1 (title + content) for content slides

    # Add a title slide as the first slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # Layout 0 = title slide
    title_slide.shapes.title.text = topic  # Set the main title to the research topic
    title_slide.placeholders[1].text = "AI Research Assistant"  # Set the subtitle

    slides = _parse_ppt_outline(writer_text)  # Parse the writer output into slide data

    if not slides:
        # Fallback: if no slide headings were detected, split by paragraph and use first 10
        for para in writer_text.split("\n\n")[:10]:
            if para.strip():
                slide = prs.slides.add_slide(blank_layout)  # Add a new content slide
                slide.shapes.title.text = para[:60]          # Use first 60 chars as title
                slide.placeholders[1].text = para            # Use full paragraph as content
    else:
        # Normal path: create one slide per parsed slide dict
        for s in slides:
            slide = prs.slides.add_slide(blank_layout)       # Add a new content slide
            slide.shapes.title.text = s["title"][:100]       # Set slide title (max 100 chars)
            tf = slide.placeholders[1].text_frame            # Get the text frame for bullet content
            tf.text = ""                                      # Clear default placeholder text
            for bullet in s["bullets"][:8]:                  # Add up to 8 bullets per slide
                p = tf.add_paragraph()                        # Add a new paragraph (bullet)
                p.text = bullet[:200]                         # Limit bullet text to 200 chars
                p.level = 0                                   # Set bullet indent level to 0 (top level)

    # Build the output file path using the topic name (spaces replaced with underscores)
    path = output_dir / f"{topic[:40].replace(' ', '_')}.pptx"
    prs.save(path)  # Save the presentation to disk
    return path     # Return the saved file path


def export_markdown(content: str, topic: str, output_dir: Path) -> Path:
    # Exports the reviewed report as a Markdown (.md) file
    path = output_dir / f"{topic[:40].replace(' ', '_')}.md"
    path.write_text(f"# {topic}\n\n{content}", encoding="utf-8")  # Add topic as H1 heading
    return path  # Return the saved file path


def export_docx(content: str, topic: str, output_dir: Path) -> Path:
    # Exports the reviewed report as a Word (.docx) file
    from docx import Document  # Import here to avoid loading if docx export is not requested
    doc = Document()  # Create a new blank Word document
    doc.add_heading(topic, 0)  # Add the topic as the document title (heading level 0)
    for para in content.split("\n\n"):  # Split content into paragraphs by double newline
        if para.strip():
            doc.add_paragraph(para.strip())  # Add each non-empty paragraph to the document
    path = output_dir / f"{topic[:40].replace(' ', '_')}.docx"
    doc.save(path)  # Save the Word document to disk
    return path     # Return the saved file path


def export_pdf(content: str, topic: str, output_dir: Path) -> Path:
    # Exports the reviewed report as a PDF file using fpdf2
    from fpdf import FPDF  # Import here to avoid loading if pdf export is not requested
    pdf = FPDF()             # Create a new blank PDF document
    pdf.add_page()           # Add the first page
    pdf.set_font("Helvetica", "B", 16)          # Set bold 16pt font for the title
    pdf.cell(0, 10, topic[:80], ln=True)        # Write the topic as the title (max 80 chars)
    pdf.set_font("Helvetica", size=11)          # Switch to regular 11pt font for body text
    for line in content.splitlines():
        # Write each line; use multi_cell to handle long lines with word wrap
        pdf.multi_cell(0, 8, line[:200] if line.strip() else "")
    path = output_dir / f"{topic[:40].replace(' ', '_')}.pdf"
    pdf.output(str(path))  # Save the PDF to disk (fpdf2 requires a string path)
    return path            # Return the saved file path


def export_all(writer_text: str, reviewer_text: str, topic: str, formats: list[str], output_dir: Path) -> dict:
    # Master export function — calls the appropriate export function for each requested format
    output_dir.mkdir(exist_ok=True)  # Ensure the output directory exists before writing files
    exported = {}  # Dictionary to collect {format: file_path} for each exported file

    if "pptx" in formats:
        exported["pptx"] = export_ppt(writer_text, topic, output_dir)  # PPT uses writer output (has slide outline)

    if "md" in formats:
        exported["md"] = export_markdown(reviewer_text, topic, output_dir)  # MD uses reviewed output

    if "docx" in formats:
        exported["docx"] = export_docx(reviewer_text, topic, output_dir)  # DOCX uses reviewed output

    if "pdf" in formats:
        exported["pdf"] = export_pdf(reviewer_text, topic, output_dir)  # PDF uses reviewed output

    return exported  # Return the dict of format → file path
